import streamlit as st
import google.generativeai as genai
from pptx.util import Pt 
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path   
import os
import re
from io import BytesIO


# ========== CONFIGURE GEMINI ==========
# 💡 Replace this with your actual Gemini API key from https://aistudio.google.com/
genai.configure(api_key="AIzaSyBHQ9w_OJDvTAyI9PerSNeFoWTtGqG0_fg")

# Use a supported model (gemini-1.5-flash or gemini-1.5-pro depending on access)
MODEL_NAME = "models/gemini-2.5-flash"


st.set_page_config(page_title="AI PPT Generator", layout="centered")
st.title("🎯 AI Presentation Generator")
st.markdown("Generate professional PowerPoint slide content using **Gemini AI**!")

# Initialize Gemini Client (assumes GOOGLE_API_KEY environment variable is set)
try:
    # Use gemini-2.5-flash for fast and reliable content generation
    model = genai.GenerativeModel(MODEL_NAME)
except Exception as e:
    st.error(f"❌ Gemini Client Initialization Failed. Please check your API key setup: {e}")
    st.stop()


# --- User Input Widgets ---

topic = st.text_input("Enter your topic:")
presentation_type = st.selectbox("Select presentation type:", ["Academic", "Corporate", "Research"])
detail_level = st.selectbox("Select detail level:", ["High-level", "Detailed"])
slide_count = st.slider("Number of slides:", 3, 15, 5)


# --- Generation Logic ---

def generate_and_download_ppt(topic, presentation_type, detail_level, slide_count):
    """Generates content via Gemini and creates a PPTX file."""

    # 1. Construct the Prompt for Gemini
    prompt = f"""
    You are an expert presentation outline creator.
    Create a {presentation_type.lower()} PowerPoint presentation on the topic "{topic}".
    It must be {detail_level.lower()} and contain exactly {slide_count} slides.
    
    Each slide must include:
    - A clear, concise title.
    - 3–5 professional bullet points.

    Return the output in the following precise format, ensuring the title is on the same line as the slide number:
    
    Slide 1: <Your Slide Title Here>
    - First Key Point
    - Second Key Point
    - Third Key Point
    ...
    
    Slide 2: <Your Slide Title Here>
    - First Key Point
    ...
    
    ... continue for exactly {slide_count} slides.
    """

    try:
        # 2. Generate content using Gemini
        response = model.generate_content(prompt)
        content = response.text
        
        # 3. Create PowerPoint presentation and parse content
        prs = Presentation()
        # Using layout 1: Title and Content
        layout = prs.slide_layouts[1] 
        
        # Split by "Slide " but START from the second element [1:]
        # This fixes the issue of the first element being empty or preamble.
        slides = content.split("Slide ")
        
        successful_slides = 0

        for s_block in slides[1:]:
            s_block = s_block.strip()
            if not s_block:
                continue

            # Find the end of the title line (first newline)
            title_line_end = s_block.find('\n')
            if title_line_end == -1:
                title_line = s_block
                body_content = ""
            else:
                title_line = s_block[:title_line_end].strip()
                body_content = s_block[title_line_end:].strip()
                
            # Use regex to robustly find content after "N:" or "N."
            match = re.match(r'^\s*\d+[:\.]\s*(.*)', title_line, re.IGNORECASE)
            if match:
                # This captures the text after the number and separator
                title = match.group(1).strip()
            else:
                # Fallback: if format is completely unexpected, skip.
                st.warning(f"Could not parse title from line: '{title_line}'. Skipping block.")
                continue

            # Extract bullet points, filtering for common bullet characters
            points = [
                line.strip("-•* ").strip() 
                for line in body_content.split("\n") 
                if line.strip().startswith(('-', '•', '*')) and line.strip()
            ]

            if not title or not points:
                st.info(f"Skipping block due to missing title or points. Title: '{title}'. Points found: {len(points)}")
                continue

            # --- Create a slide ---
            slide = prs.slides.add_slide(layout)

            # Title Formatting
            title_shape = slide.shapes.title
            title_shape.text = title
            # Example custom formatting (optional, but good practice)
            title_shape.text_frame.paragraphs[0].font.name = "Arial"
            title_shape.text_frame.paragraphs[0].font.size = Pt(28)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

            # Content Formatting
            content_box = slide.placeholders[1]
            content_box.text = ""  # clear default text
            
            for point in points:
                p = content_box.text_frame.add_paragraph()
                p.text = point
                p.level = 0 # Ensures it's a top-level bullet
                p.font.name = "Arial"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(51, 51, 51) # Gray
                
            successful_slides += 1

        if successful_slides == 0:
            st.error("Could not generate or parse any complete slides. Please refine your topic.")
            st.code(f"Raw API Content:\n{content}")
            return None, None
            
        # 4. Save presentation to an in-memory buffer (BytesIO)
        # This is necessary for Streamlit download button
        binary_output = BytesIO()
        prs.save(binary_output)
        
        # Rewind the buffer's position to the start before reading/downloading
        binary_output.seek(0)
        
        filename = f"{topic.replace(' ', '_')}_Presentation.pptx"
        return binary_output.getvalue(), filename

    except Exception as e:
        st.error(f"❌ An error occurred during presentation generation: {str(e)}")
        return None, None


# --- Button Click Handler ---

if st.button("Generate Presentation"):
    if not topic.strip():
        st.warning("⚠️ Please enter a topic first.")
    else:
        with st.spinner(f"✨ Generating content and compiling {slide_count} slides..."):
            ppt_data, filename = generate_and_download_ppt(topic, presentation_type, detail_level, slide_count)

            if ppt_data and filename:
                st.success(f"✅ Presentation compiled successfully! ({len(ppt_data)/1024:.2f} KB)")
                
                # Display the download button
                st.download_button(
                    label="⬇️ Download PPTX File",
                    data=ppt_data,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key='download_ppt_button'
                )

                st.markdown("---")
                st.info("💡 **Next Steps:** Download and open the file. You can then enhance the visual design in PowerPoint!")
